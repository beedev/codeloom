"""Generate CodeLoom Architecture PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand colors ──────────────────────────────────────────────────
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)       # Deep navy
ACCENT_BLUE = RGBColor(0x38, 0x8B, 0xFD)   # Electric blue
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xAA)   # Teal/cyan
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6) # Purple
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x42)  # Orange
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xC8)
MED_GRAY = RGBColor(0x6B, 0x72, 0x80)
CARD_BG = RGBColor(0x1A, 0x23, 0x3B)       # Slightly lighter navy

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_card(slide, left, top, width, height, title, items, accent=ACCENT_BLUE):
    # Card background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                  Inches(0.06), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    # Title
    add_text(slide, left + 0.2, top + 0.1, width - 0.3, 0.4, title,
             font_size=16, color=accent, bold=True)
    # Items
    y = top + 0.5
    for item in items:
        add_text(slide, left + 0.2, y, width - 0.4, 0.3, item,
                 font_size=11, color=LIGHT_GRAY)
        y += 0.28
    return shape


def add_stat_box(slide, left, top, number, label, accent=ACCENT_CYAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                    Inches(1.8), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()
    add_text(slide, left + 0.1, top + 0.08, 1.6, 0.5, number,
             font_size=28, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left + 0.1, top + 0.65, 1.6, 0.4, label,
             font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_dark_bg(slide)

add_text(slide, 1.0, 1.5, 11, 1.2, "CodeLoom", font_size=54, color=ACCENT_BLUE, bold=True)
add_text(slide, 1.0, 2.7, 11, 0.8, "AI-Powered Code Intelligence & Migration Platform",
         font_size=28, color=WHITE)
add_text(slide, 1.0, 3.6, 11, 0.6,
         "AST + ASG + RAG + Multi-Lane Migration Engine + Full-Autonomy CLI Orchestration",
         font_size=16, color=LIGHT_GRAY)

# Accent line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.4),
                               Inches(4), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_CYAN
line.line.fill.background()

add_text(slide, 1.0, 5.5, 5, 0.4, "Architecture & Capabilities Overview",
         font_size=14, color=MED_GRAY)
add_text(slide, 1.0, 5.9, 5, 0.4, "March 2026",
         font_size=12, color=MED_GRAY)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: Platform Overview
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)

add_text(slide, 0.8, 0.3, 12, 0.6, "Platform Overview", font_size=32, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.85),
                               Inches(2.5), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

add_text(slide, 0.8, 1.1, 11, 0.5,
         "CodeLoom is a two-layer system: a code intelligence platform + an LLM-powered migration orchestrator.",
         font_size=14, color=LIGHT_GRAY)

# Stats row
add_stat_box(slide, 0.8, 1.7, "14", "Languages\nParsed", ACCENT_BLUE)
add_stat_box(slide, 2.8, 1.7, "9", "ASG Edge\nTypes", ACCENT_CYAN)
add_stat_box(slide, 4.8, 1.7, "20", "MCP\nTools", ACCENT_PURPLE)
add_stat_box(slide, 6.8, 1.7, "4", "Migration\nLanes", ACCENT_ORANGE)
add_stat_box(slide, 8.8, 1.7, "6", "Skill\nPhases", ACCENT_BLUE)
add_stat_box(slide, 10.8, 1.7, "96%", "Proven\nAccuracy", ACCENT_CYAN)

# Architecture layers
add_card(slide, 0.8, 3.2, 5.8, 4.0, "Layer 1: CodeLoom Platform", [
    "FastAPI + PostgreSQL + pgvector backend",
    "AST Parsers: tree-sitter (Python/JS/TS/Java/C#/COBOL)",
    "  + regex (PL/I, JCL, VB.NET, SQL, JSP, ASP, XML)",
    "ASG Builder: 9 edge types incl. COBOL PERFORM/CALL",
    "Migration Engine: V1/V2/V3 pipelines + agentic loop",
    "4 Migration Lanes with source type context",
    "RAG Pipeline: BM25 + vector + RAPTOR + reranking",
    "Knowledge Projects: composable document notebooks",
    "Understanding Engine: deep analysis + entry points",
    "Reverse Engineering: 9-chapter auto-documentation",
    "20 MCP tools for external orchestration",
], accent=ACCENT_BLUE)

add_card(slide, 6.9, 3.2, 5.8, 4.0, "Layer 2: /migrate Skill (Claude Code CLI)", [
    "6 phase skills: init / run / compare / refactor / resume / status",
    "5 lane sub-skills (mainframe, mainframe-transform, etc.)",
    "30-rule learned system (gate-keyed, auto-demotion)",
    "Source-first accuracy methodology (CHECK A-F)",
    "Cross-MVP coordination: SYMBOLS.md + SPEC.md",
    "Checkpoint/resume (survives /clear)",
    "Migration Manifest: Detect -> Propose -> Confirm",
    "Knowledge enrichment: project + reference notebooks",
    "Compile gate per MVP (3 fix rounds max)",
    "Only 1 user approval gate (init Step 6)",
    "Learning capture: LESSONS.md + rules.md updates",
], accent=ACCENT_PURPLE)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3: AST Parsing & ASG
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)

add_text(slide, 0.8, 0.3, 12, 0.6, "AST Parsing & Abstract Semantic Graph",
         font_size=32, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.85),
                               Inches(4), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_CYAN
line.line.fill.background()

# Tree-sitter parsers
add_card(slide, 0.8, 1.2, 3.8, 2.8, "Tree-Sitter Parsers (High Fidelity)", [
    "Python (.py) -- functions, classes, methods",
    "JavaScript/TypeScript (.js/.ts/.jsx/.tsx)",
    "Java (.java) -- with optional JavaParser bridge",
    "C# (.cs) -- with optional Roslyn bridge",
    "COBOL (.cbl/.cob/.cpy) -- programs, paragraphs,",
    "  sections, copybooks, EXEC SQL/CICS/DLI detection",
], accent=ACCENT_BLUE)

# Regex parsers
add_card(slide, 4.8, 1.2, 3.8, 2.8, "Regex Parsers", [
    "PL/I (.pl1/.pli) -- procedures, entries, packages",
    "JCL (.jcl/.proc) -- jobs, steps, procs, DD stmts",
    "VB.NET (.vb) -- classes, modules, subs, functions",
    "SQL (.sql) -- stored procs, views, triggers",
    "JSP (.jsp) -- Java Server Pages",
    "ASP.NET (.aspx/.ascx) -- Web Forms",
    "XML Config, Properties files",
], accent=ACCENT_CYAN)

# ASG edge types
add_card(slide, 8.8, 1.2, 3.8, 2.8, "9 ASG Edge Types", [
    "contains -- class -> method nesting",
    "imports -- file -> file (COPY, import, using)",
    "calls -- PERFORM, CALL, function invocation",
    "inherits -- class -> base class",
    "implements -- class -> interface",
    "overrides -- method -> parent method",
    "type_dep -- field/param type references",
    "calls_sp -- app code -> stored procedure",
    "data_flow -- JCL step -> step (dataset sharing)",
], accent=ACCENT_PURPLE)

# COBOL-specific detail
add_card(slide, 0.8, 4.3, 5.8, 2.8, "COBOL Parser Deep Dive (tree-sitter)", [
    "Program, paragraph, section, copybook unit extraction",
    "COPY statements -> imports edges for dependency graph",
    "PERFORM / CALL / GO TO -> calls edges (with THRU support)",
    "EXEC SQL/CICS/DLI pre-stripped, restored post-parse, metadata flags",
    "IMS DL/I function code extraction (GU, GN, ISRT, REPL, DLET...)",
    "Program classification: batch_program, cics_online, utility_subprogram, ims_dli",
    "SELECT...ASSIGN TO -> file-DDNAME mapping extraction",
    "Error node recovery: regex fallback for orphaned paragraphs",
], accent=ACCENT_ORANGE)

# JCL detail
add_card(slide, 6.9, 4.3, 5.8, 2.8, "JCL Parser + Data Flow Edges", [
    "Fixed-format column parsing (cols 1-72)",
    "JOB, EXEC PGM, EXEC PROC, inline PROC...PEND extraction",
    "DD statement metadata: DSN, DISP, DDNAME per step",
    "Step classification: compile_link, sort_merge, data_mgmt, application_run",
    "EXEC PGM -> calls edge to COBOL/PL/I programs",
    "data_flow edges: DISP=NEW/PASS (producer) -> DISP=SHR/OLD (consumer)",
    "Temporary dataset (&&name) flow tracked within same job",
    "PL/I parser: CALL/GO TO edges, %INCLUDE imports, ON conditions",
], accent=ACCENT_BLUE)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4: Migration Engine & Lanes
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)

add_text(slide, 0.8, 0.3, 12, 0.6, "Migration Engine & Lanes",
         font_size=32, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.85),
                               Inches(3), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_ORANGE
line.line.fill.background()

# Pipeline versions
add_card(slide, 0.8, 1.2, 3.8, 2.5, "Pipeline Versions", [
    "V1 (Legacy): 6-phase Discovery-first",
    "V2 (Default): 4-phase Architecture-first",
    "  Architecture -> Discovery -> Transform -> Test",
    "V3: 5-phase adds auto-flowing Design phase",
    "Batch execution: sequential MVPs, auto-approval",
    "Checkpoint/resume for failed transforms",
    "Max 3 retries with exponential backoff",
], accent=ACCENT_BLUE)

# Agentic loop
add_card(slide, 4.8, 1.2, 3.8, 2.5, "Agentic Migration Loop", [
    "Multi-turn LLM conversation with tool-calling",
    "10 agent tools: read source, search code,",
    "  lookup docs, validate syntax",
    "Max turns configurable (default 10)",
    "Streaming AgentEvents for SSE to frontend",
    "Lane augment_prompt() injects domain knowledge",
    "Ground truth validation (advisory)",
], accent=ACCENT_CYAN)

# Quality gates
add_card(slide, 8.8, 1.2, 3.8, 2.5, "Quality Gates", [
    "6 gate categories: PARITY, COMPILE, UNIT_TEST,",
    "  INTEGRATION, CONTRACT, REGRESSION",
    "Blocking vs advisory gates per lane",
    "Per-rule confidence scoring (0.0-1.0)",
    "Confidence tiers: high (>=0.90), standard,low",
    "Stub quality check (TODO/FIXME detection)",
    "Compile gate: hard stop after 3 failed rounds",
], accent=ACCENT_PURPLE)

# 4 lanes
lanes = [
    ("Mainframe -> Modern", ACCENT_ORANGE, [
        "COBOL/PL/I/JCL/CICS/IMS -> Python/Java/.NET",
        "7 deterministic transform rules",
        "16 source type semantics for LLM proposals",
        "COMP-3->Decimal, STOP RUN->SystemExit, VSAM->KV",
    ]),
    ("Struts -> Spring Boot", ACCENT_BLUE, [
        "Struts 1.x/2.x -> Spring Boot REST",
        "Action->@Controller, ActionForm->@RequestBody",
        "10 source type semantics (JSP, tiles, validation)",
        "struts-config.xml -> annotations",
    ]),
    ("StoredProc -> ORM", ACCENT_CYAN, [
        "SQL SP/View/Trigger -> JPA/SQLAlchemy/EF",
        "CRUD SP -> repository method",
        "8 source type semantics",
        "DDL schema -> ORM entity definitions",
    ]),
    ("VB.NET -> .NET Core", ACCENT_PURPLE, [
        "VB.NET Web Forms -> C#/ASP.NET Core MVC",
        "WebForm->Razor, Module->static class",
        "10 source type semantics",
        "Web.config -> appsettings.json",
    ]),
]

x_pos = 0.8
for title, accent, items in lanes:
    add_card(slide, x_pos, 4.1, 2.95, 3.0, title, items, accent=accent)
    x_pos += 3.1

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5: /migrate Skill Phases
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)

add_text(slide, 0.8, 0.3, 12, 0.6, "/migrate: Full-Autonomy Migration Skill",
         font_size=32, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.85),
                               Inches(4), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_PURPLE
line.line.fill.background()

add_text(slide, 0.8, 1.0, 11, 0.5,
         "Claude Code CLI is the orchestrator. CodeLoom is the read-only intelligence source (AST, ASG, code units, RAG).",
         font_size=13, color=LIGHT_GRAY)

phases = [
    ("1. Init", ACCENT_BLUE, [
        "Project selection + intel gathering",
        "Reverse engineering context",
        "Lane auto-detection (confidence scoring)",
        "Knowledge notebook selection",
        "Architecture doc + Migration Manifest",
        "MVP clustering + SPEC.md per MVP",
        "Uncovered file audit",
        "User approval gate (ONLY gate)",
    ]),
    ("2. Run", ACCENT_CYAN, [
        "Per-MVP batch transform loop",
        "Source units pulled via MCP",
        "Knowledge enrichment (Step 3.5)",
        "Migration manifest drives targets",
        "SYMBOLS.md cross-MVP imports",
        "Compile gate (3 rounds max)",
        "Validate + complete via MCP",
        "Chains to next MVP automatically",
    ]),
    ("3. Compare", ACCENT_ORANGE, [
        "Source-first methodology (NON-NEGOTIABLE)",
        "Pull ALL source units via MCP",
        "CHECK A-F per construct:",
        "  Presence, Branching, Boundaries,",
        "  Calls, Data, Returns",
        "Score: (correct+deviation) / total",
        "Auto-fix surgical bugs (<= 3 lines)",
        "Re-migration trigger < 40%",
    ]),
    ("4. Refactor", ACCENT_PURPLE, [
        "Idiomatic target patterns",
        "Context managers, match/case",
        "Dataclasses, type hints, generators",
        "Extract services, custom exceptions",
        "Pass 2 accuracy comparison",
        "Regression check (>5pt drop -> revert)",
        "No functional behavior changes",
        "codeloom_save_accuracy_report",
    ]),
]

x_pos = 0.8
for title, accent, items in phases:
    add_card(slide, x_pos, 1.5, 2.95, 5.5, title, items, accent=accent)
    x_pos += 3.1

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6: Knowledge System & Migration Manifest
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)

add_text(slide, 0.8, 0.3, 12, 0.6, "Knowledge System & Migration Manifest",
         font_size=32, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.85),
                               Inches(4), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_CYAN
line.line.fill.background()

# Knowledge system
add_card(slide, 0.8, 1.2, 5.8, 3.0, "Composable Knowledge Notebooks", [
    "Upload any document: IBM Redbooks, client SOW, architecture standards",
    "Chunk -> embed -> pgvector with BM25 + vector + RAPTOR + reranking",
    "Project notebooks: client-specific (business rules, data models, compliance)",
    "Reference notebooks: shared (PL/I Manual, CICS Redbook, AWS docs)",
    "Queried during migration: init (architecture), run (transform), compare (accuracy)",
    "codeloom_search_knowledge(notebook_id, query) MCP tool",
    "BM25 handles exact IBM terms (DFHCOMMAREA, SQLCA) precisely",
    "More flexible than fixed catalogs -- each migration assembles its own context",
], accent=ACCENT_CYAN)

# Migration manifest
add_card(slide, 6.9, 1.2, 5.8, 3.0, "Migration Manifest (NEW)", [
    "Platform DETECTS source types (parsers + ASG -- deterministic)",
    "LLM PROPOSES targets (intelligent, context-aware via lane semantics)",
    "User CONFIRMS (final authority at approval gate)",
    "",
    "Lane provides source type SEMANTICS, not hardcoded maps",
    "LLM reasons about the right target per project context",
    "Unknown source types flagged for user input",
    "Persisted as target_manifest JSONB on MigrationPlan",
], accent=ACCENT_ORANGE)

# Manifest flow diagram
add_card(slide, 0.8, 4.5, 11.8, 2.7, "Detect -> Propose -> Confirm Flow", [
    "DETECT: COBOL parser classifies programs (batch/CICS/utility/IMS). JCL parser classifies steps (sort/compile/run).",
    "         get_source_type_inventory() aggregates: {source_type, count, sample_units, flags}",
    "",
    "PROPOSE: Lane's get_source_type_context() describes what each type IS (semantics).",
    "         LLM uses inventory + lane context + user's target ecosystem to propose targets with reasoning.",
    "",
    "CONFIRM: Approval gate shows: Auto (high confidence) | Review (medium) | User choice (configurable).",
    "         User overrides as needed. Finalized manifest saved to plan and flows to SPEC.md per MVP.",
], accent=ACCENT_PURPLE)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7: MCP Tools & Integration
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)

add_text(slide, 0.8, 0.3, 12, 0.6, "20 MCP Tools: Platform as API",
         font_size=32, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.85),
                               Inches(3), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

add_text(slide, 0.8, 1.0, 11, 0.5,
         "CodeLoom exposes its intelligence as 20 MCP tools -- any LLM orchestrator (Claude Code, agents, CI/CD) can consume them.",
         font_size=13, color=LIGHT_GRAY)

tool_groups = [
    ("Project Intelligence", ACCENT_BLUE, [
        "codeloom_list_projects",
        "codeloom_get_project_intel",
        "  + source_type_inventory (NEW)",
    ]),
    ("Source Code Access", ACCENT_CYAN, [
        "codeloom_list_units (paginated, filtered)",
        "codeloom_get_source_unit (full source + metadata)",
        "codeloom_search_codebase (RAG semantic search)",
        "codeloom_search_knowledge (notebook RAG)",
    ]),
    ("Migration Planning", ACCENT_PURPLE, [
        "codeloom_get_lane_info + source_type_context",
        "codeloom_get_import_graph (fan-in analysis)",
        "codeloom_save_plan + target_manifest (NEW)",
        "codeloom_save_mvps (file path resolution)",
    ]),
    ("Migration Execution", ACCENT_ORANGE, [
        "codeloom_start_transform",
        "codeloom_complete_transform",
        "codeloom_validate_output (ground truth)",
        "codeloom_save_accuracy_report",
    ]),
    ("MVP Context", ACCENT_CYAN, [
        "codeloom_list_mvps (status, confidence)",
        "codeloom_get_mvp_context (units, lane, analysis)",
        "codeloom_get_compiled_context (dep-ordered)",
    ]),
    ("Documentation", ACCENT_BLUE, [
        "codeloom_generate_reverse_doc (9-chapter)",
        "codeloom_get_reverse_doc",
        "codeloom_list_reverse_docs",
    ]),
]

x_pos = 0.8
y_pos = 1.5
for i, (title, accent, items) in enumerate(tool_groups):
    add_card(slide, x_pos, y_pos, 3.8, 2.2, title, items, accent=accent)
    x_pos += 4.0
    if i == 2:
        x_pos = 0.8
        y_pos = 3.9

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8: Accuracy Methodology
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)

add_text(slide, 0.8, 0.3, 12, 0.6, "Source-First Accuracy Methodology",
         font_size=32, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.85),
                               Inches(4), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_ORANGE
line.line.fill.background()

add_text(slide, 0.8, 1.0, 11, 0.4,
         "NEVER score accuracy by reading only target files. Pull ALL source, count every construct, find each target equivalent.",
         font_size=14, color=ACCENT_ORANGE, bold=True)

# 6 checks
checks = [
    ("CHECK A", "Presence", "Target function covers this source construct?"),
    ("CHECK B", "Branching", "All IF/EVALUATE/ON branches represented?"),
    ("CHECK C", "Boundaries", "Comparison operators correct? (> vs >=)"),
    ("CHECK D", "Calls", "CALL/PERFORM -> correct target function calls?"),
    ("CHECK E", "Data", "Field/variable references -> correct attributes?"),
    ("CHECK F", "Returns", "RETURN-CODE / output fields -> exceptions or return values?"),
]

x_pos = 0.8
for code, name, desc in checks:
    add_card(slide, x_pos, 1.6, 1.95, 1.8, f"{code}: {name}", [desc], accent=ACCENT_CYAN)
    x_pos += 2.05

# Classification
add_card(slide, 0.8, 3.7, 5.8, 1.5, "Classification", [
    "Correct -- logic match confirmed",
    "Gap -- no equivalent in target (conservative: if unsure, classify as Gap)",
    "Bug -- wrong logic (wrong operator, missing branch)",
    "Deviation -- intentional architectural change (documented reason)",
], accent=ACCENT_BLUE)

# Scoring
add_card(slide, 6.9, 3.7, 5.8, 1.5, "Scoring & Gates", [
    "Weights: main paragraphs x3, entry points x2, utility x1",
    "Score = (correct + deviation) / total_weighted x 100",
    "Re-migration trigger: < 40% -> re-generate from scratch",
    "Line ratio red flag: target/source < 0.15 -> investigate",
], accent=ACCENT_ORANGE)

# Proven results
add_card(slide, 0.8, 5.5, 11.8, 1.7, "Proven Results: lema01 Project (37 COBOL + JCL files)", [
    "96/100 accuracy -- 86 source constructs, 79 correct, 7 justified deviations, 0 bugs, 0 gaps",
    "Every comparison operator verified (> vs >= across 5 programs with 3 different semantics)",
    "VSS NS-path REWRITE anomaly preserved (QOH decremented even on insufficient stock)",
    "85/85 tests passing across all 4 MVPs (Foundation 16, Order Processing 31, Billing & VSAM 51, Utilities 85)",
], accent=ACCENT_PURPLE)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9: Continuous Learning
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)

add_text(slide, 0.8, 0.3, 12, 0.6, "Continuous Learning & Reverse Engineering",
         font_size=32, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.85),
                               Inches(4), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_PURPLE
line.line.fill.background()

add_card(slide, 0.8, 1.2, 5.8, 3.2, "Learned Rules System", [
    "30-rule max, gate-keyed (init/arch, run/gen, run/compile, compare)",
    "Hit count tracking -- rules promoted after 2+ occurrences",
    "Auto-demotion after 5 migrations without trigger",
    "LESSONS.md per migration: systemic vs one-off tagged",
    "",
    "Current rules (11 active):",
    "  R1: PK strategy from source data model",
    "  R8: Source-first scoring (2 hits -- prevented hallucinated scores)",
    "  R10: COBOL PIC width verification in test data",
    "  R11: Per-program comparison operator verification",
], accent=ACCENT_PURPLE)

add_card(slide, 6.9, 1.2, 5.8, 3.2, "Reverse Engineering Documentation", [
    "9-chapter auto-generated docs from source code",
    "  Ch 1: System Overview (deterministic from ASG)",
    "  Ch 2: Architecture Analysis (LLM-synthesized)",
    "  Ch 3: Component Inventory (from parser units)",
    "  Ch 4: Data Flow Analysis (from ASG edges)",
    "  Ch 5: Business Logic Extraction (LLM-synthesized)",
    "  Ch 6: Integration Points (from import graph)",
    "  Ch 7: Technology Stack (from file breakdown)",
    "  Ch 8: External Dependencies (from call analysis)",
    "  Ch 9: Error Handling Patterns (from code search)",
], accent=ACCENT_CYAN)

# Understanding engine
add_card(slide, 0.8, 4.7, 5.8, 2.5, "Understanding Engine (Deep Analysis)", [
    "Tiered token budgets: Tier 1 (<100K full), Tier 2 (<200K truncated), Tier 3 (summarized)",
    "Entry point detection: zero-incoming-calls + annotation patterns",
    "Call chain tracing: transitive closure of calls edges",
    "Business rule extraction: domain logic identification",
    "Background worker: poll 15s, max 2 concurrent, heartbeat + stale reclamation",
    "Coverage metrics: warn < 50%, target 80%",
], accent=ACCENT_ORANGE)

# RAG pipeline
add_card(slide, 6.9, 4.7, 5.8, 2.5, "RAG Pipeline", [
    "BM25 lexical ranking (exact term match for IBM components)",
    "Vector similarity search (semantic, pgvector HNSW index)",
    "Hybrid retrieval: BM25 + vector combined scoring",
    "RAPTOR: hierarchical tree-based summaries for broad context",
    "Reranking: mixedbread-ai/mxbai-rerank-base-v1",
    "Stateless query: thread-safe, multi-user (not global state mutation)",
], accent=ACCENT_BLUE)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10: Competitive Positioning
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)

add_text(slide, 0.8, 0.3, 12, 0.6, "Competitive Advantages",
         font_size=32, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.85),
                               Inches(3), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_CYAN
line.line.fill.background()

advantages = [
    ("Auditable Accuracy", ACCENT_BLUE,
     "Source-first, paragraph-by-paragraph comparison with CHECK A-F.\n"
     "Not self-reported -- verifiable methodology. 96% proven on real project."),
    ("Composable Knowledge", ACCENT_CYAN,
     "Notebooks absorb IBM Redbooks + client SOW + target standards +\n"
     "compliance docs. Each migration assembles its own context."),
    ("Structural Intelligence", ACCENT_PURPLE,
     "tree-sitter COBOL parser + 9 ASG edge types + EXEC SQL/CICS/IMS\n"
     "metadata. Programmatic access to source structure, not just text."),
    ("Post-Migration Value", ACCENT_ORANGE,
     "After migration, CodeLoom remains a code intelligence platform:\n"
     "ASG browsing, RAG chat, understanding engine, diagrams."),
    ("Operator-Level Precision", ACCENT_BLUE,
     "Obsessive focus on > vs >=, VSAM anomalies, PIC widths, COMP-3.\n"
     "Where migrations actually succeed or fail in production."),
    ("Human-in-the-Loop", ACCENT_CYAN,
     "The /migrate skill lets an engineer steer at every phase.\n"
     "For enterprise migrations, this is a feature, not a limitation."),
]

x_pos = 0.8
y_pos = 1.2
for i, (title, accent, desc) in enumerate(advantages):
    add_card(slide, x_pos, y_pos, 3.8, 1.8, title, desc.split("\n"), accent=accent)
    x_pos += 4.0
    if i == 2:
        x_pos = 0.8
        y_pos = 3.3

# Bottom tagline
add_text(slide, 0.8, 5.5, 11.8, 0.8,
         "CodeLoom: Code intelligence platform that understands your codebase structurally,\n"
         "migrates it accurately, and remains valuable after migration is complete.",
         font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════
output_path = "/Users/bharath/Desktop/codeloom/CodeLoom-Architecture.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Slides: {len(prs.slides)}")
